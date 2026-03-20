#!/usr/bin/env python3
"""Spike 004: PDF/PPTX Conversion Quality Assessment.

Tests MarkItDown conversion of PDF, PPTX, and DOCX documents.
Evaluates text extraction, image extraction, table fidelity,
speaker notes inclusion, and compares alternative extraction
approaches (PyMuPDF for PDF images, python-pptx for speaker notes).

Usage:
    cd src/functions && uv run python ../spikes/004-pdf-pptx-conversion/run.py

Prerequisites (install into functions venv):
    uv pip install pymupdf python-pptx python-docx reportlab Pillow
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------
SPIKE_DIR = Path(__file__).resolve().parent
SAMPLES_DIR = SPIKE_DIR / "samples"
OUTPUT_DIR = SPIKE_DIR / "output"

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def count_headings(md: str) -> int:
    return len(re.findall(r"^#{1,6}\s+", md, re.MULTILINE))


def count_links(md: str) -> int:
    return len(re.findall(r"\[([^\]]+)\]\(([^)]+)\)", md))


def extract_links(md: str) -> list[tuple[str, str]]:
    """Extract Markdown hyperlinks as (text, url) tuples.

    Excludes image references (![alt](src)) — only counts [text](url).
    """
    # (?<!!) = negative lookbehind to exclude ![alt](url) image references
    return re.findall(r"(?<!!)\[([^\]]+)\]\((https?://[^)]+)\)", md)


def extract_urls(md: str) -> list[str]:
    """Extract all bare URLs (http/https) from markdown text."""
    return re.findall(r"https?://[^\s\)>\"]+", md)


def count_tables(md: str) -> int:
    return len(re.findall(r"^\|.+\|$", md, re.MULTILINE))


def count_images_md(md: str) -> int:
    """Count image references in markdown (![alt](src) pattern)."""
    return len(re.findall(r"!\[([^\]]*)\]\(([^)]+)\)", md))


def section_header(title: str) -> None:
    print(f"\n{'=' * 70}")
    print(f"  {title}")
    print(f"{'=' * 70}")


def sub_header(title: str) -> None:
    print(f"\n--- {title} ---")


# ---------------------------------------------------------------------------
# 1. PDF Conversion Tests
# ---------------------------------------------------------------------------

def test_pdf_markitdown(pdf_path: Path) -> str:
    """Test MarkItDown PDF conversion and return Markdown output."""
    from markitdown import MarkItDown

    sub_header("MarkItDown PDF Conversion")
    md = MarkItDown()
    result = md.convert(str(pdf_path))
    text = result.text_content
    print(f"  Characters: {len(text):,}")
    print(f"  Lines: {len(text.splitlines()):,}")
    print(f"  Headings: {count_headings(text)}")
    print(f"  Links: {count_links(text)}")
    print(f"  Table rows: {count_tables(text)}")
    print(f"  Image refs: {count_images_md(text)}")

    # Check for expected content
    checks = {
        "Title present": "Azure Knowledge Base Architecture Guide" in text
                         or "Azure Knowledge Base" in text,
        "Table content": "Azure AI Search" in text or "AI Search" in text,
        "Sub-heading": "Component" in text or "Deployment" in text,
        "Bullet points": "managed identity" in text.lower() or "retrieval" in text.lower(),
    }
    for label, passed in checks.items():
        status = "✅" if passed else "❌"
        print(f"  {status} {label}")

    # Hyperlink check
    md_links = extract_links(text)
    bare_urls = extract_urls(text)
    print(f"\n  Hyperlinks (Markdown [text](url)): {len(md_links)}")
    for link_text, link_url in md_links[:5]:
        print(f"    [{link_text}]({link_url})")
    print(f"  Bare URLs found: {len(bare_urls)}")
    for url in bare_urls[:5]:
        print(f"    {url}")
    total_urls = len(md_links) + len(bare_urls)
    if total_urls > 0:
        print(f"  ✅ Hyperlinks detected ({total_urls} total)")
    else:
        print(f"  ⚠️  No hyperlinks detected in PDF conversion output")

    # Long table check
    table_rows = count_tables(text)
    print(f"\n  Table rows in output: {table_rows}")
    if table_rows >= 30:
        print(f"  ✅ Long table (30+ rows) appears in output")
    elif table_rows >= 20:
        print(f"  ⚠️  Long table partially extracted ({table_rows} rows)")
    else:
        print(f"  ⚠️  Long table may be split or truncated (only {table_rows} rows)")

    return text


def test_pdf_image_extraction_pymupdf(pdf_path: Path) -> list[dict]:
    """Extract images from PDF using PyMuPDF (fitz)."""
    import fitz  # PyMuPDF

    sub_header("PyMuPDF PDF Image Extraction")
    doc = fitz.open(str(pdf_path))
    images = []
    for page_num in range(len(doc)):
        page = doc[page_num]
        image_list = page.get_images(full=True)
        for img_index, img_info in enumerate(image_list):
            xref = img_info[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            width = base_image["width"]
            height = base_image["height"]
            colorspace = base_image.get("colorspace", "unknown")

            images.append({
                "page": page_num + 1,
                "index": img_index,
                "xref": xref,
                "ext": image_ext,
                "width": width,
                "height": height,
                "size_bytes": len(image_bytes),
                "colorspace": colorspace,
                "bytes": image_bytes,
            })
            print(f"  Page {page_num + 1}, Image {img_index}: "
                  f"{width}x{height} {image_ext} ({len(image_bytes):,} bytes) "
                  f"colorspace={colorspace}")

    doc.close()
    print(f"  Total images extracted: {len(images)}")
    return images


def test_pdf_image_extraction_markitdown(pdf_path: Path) -> str:
    """Check what MarkItDown reports about images in PDF."""
    from markitdown import MarkItDown

    sub_header("MarkItDown PDF Image Handling")
    md = MarkItDown()
    result = md.convert(str(pdf_path))
    text = result.text_content

    image_refs = re.findall(r"!\[([^\]]*)\]\(([^)]+)\)", text)
    if image_refs:
        print(f"  Found {len(image_refs)} image reference(s) in Markdown:")
        for alt, src in image_refs:
            print(f"    ![{alt}]({src})")
    else:
        print("  ⚠️  No image references found in MarkItDown PDF output")
        print("  Note: MarkItDown may not extract embedded PDF images as references")

    # Check for image-related text
    for pattern in ["image", "figure", "diagram", "chart", "png", "jpg"]:
        if pattern in text.lower():
            print(f"  Found '{pattern}' mentioned in text")

    return text


# ---------------------------------------------------------------------------
# 2. PPTX Conversion Tests
# ---------------------------------------------------------------------------

def test_pptx_markitdown(pptx_path: Path) -> str:
    """Test MarkItDown PPTX conversion."""
    from markitdown import MarkItDown

    sub_header("MarkItDown PPTX Conversion")
    md = MarkItDown()
    result = md.convert(str(pptx_path))
    text = result.text_content
    print(f"  Characters: {len(text):,}")
    print(f"  Lines: {len(text.splitlines()):,}")
    print(f"  Headings: {count_headings(text)}")
    print(f"  Links: {count_links(text)}")
    print(f"  Table rows: {count_tables(text)}")
    print(f"  Image refs: {count_images_md(text)}")

    # Check for expected content
    checks = {
        "Title present": "Azure KB Agent" in text or "Architecture" in text,
        "Slide content": "System Components" in text or "Core Services" in text,
        "Table content": "MarkItDown" in text or "Backend" in text,
        "Bullet points": "Azure AI Search" in text or "Azure OpenAI" in text,
    }
    for label, passed in checks.items():
        status = "✅" if passed else "❌"
        print(f"  {status} {label}")

    # Critical check: speaker notes
    speaker_notes_keywords = [
        "managed identity",
        "connection strings",
        "DefaultAzureCredential",
        "data flow",
        "cost-effective",
    ]
    found_notes = [kw for kw in speaker_notes_keywords if kw.lower() in text.lower()]
    if found_notes:
        print(f"  ✅ Speaker notes content detected ({len(found_notes)}/{len(speaker_notes_keywords)} keywords)")
        for kw in found_notes:
            print(f"    Found: '{kw}'")
    else:
        print(f"  ❌ Speaker notes NOT detected in MarkItDown output")
        print(f"    None of these keywords found: {speaker_notes_keywords}")

    # Hyperlink check
    md_links = extract_links(text)
    bare_urls = extract_urls(text)
    print(f"\n  Hyperlinks (Markdown [text](url)): {len(md_links)}")
    for link_text, link_url in md_links[:5]:
        print(f"    [{link_text}]({link_url})")
    print(f"  Bare URLs found: {len(bare_urls)}")
    for url in bare_urls[:5]:
        print(f"    {url}")
    total_urls = len(md_links) + len(bare_urls)
    if total_urls > 0:
        print(f"  ✅ Hyperlinks detected ({total_urls} total)")
    else:
        print(f"  ⚠️  No hyperlinks detected in PPTX conversion output")

    # Long table check
    table_rows = count_tables(text)
    print(f"\n  Table rows in output: {table_rows}")
    if table_rows >= 30:
        print(f"  ✅ Long table (30+ rows) preserved as single table")
    elif table_rows >= 20:
        print(f"  ⚠️  Long table partially extracted ({table_rows} rows)")

    return text


def test_pptx_speaker_notes_native(pptx_path: Path) -> list[dict]:
    """Extract speaker notes directly using python-pptx."""
    from pptx import Presentation

    sub_header("python-pptx Speaker Notes Extraction")
    prs = Presentation(str(pptx_path))
    slides_data = []

    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame and shape.shape_type is not None:
                if hasattr(shape, "text") and shape == slide.shapes.title:
                    title = shape.text

        notes_text = ""
        if slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_text = notes_slide.notes_text_frame.text

        slides_data.append({
            "slide": slide_num,
            "title": title,
            "notes": notes_text,
            "notes_chars": len(notes_text),
        })
        print(f"  Slide {slide_num}: '{title}'")
        if notes_text:
            preview = notes_text[:100] + ("..." if len(notes_text) > 100 else "")
            print(f"    Notes ({len(notes_text)} chars): {preview}")
        else:
            print(f"    Notes: (none)")

    total_notes = sum(s["notes_chars"] for s in slides_data)
    slides_with_notes = sum(1 for s in slides_data if s["notes_chars"] > 0)
    print(f"\n  Summary: {slides_with_notes}/{len(slides_data)} slides have speaker notes "
          f"({total_notes:,} total chars)")

    return slides_data


def test_pptx_image_extraction(pptx_path: Path) -> list[dict]:
    """Extract images from PPTX using python-pptx."""
    from pptx import Presentation

    sub_header("python-pptx Image Extraction")
    prs = Presentation(str(pptx_path))
    images = []

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                image = shape.image
                blob = image.blob
                content_type = image.content_type
                width = shape.width
                height = shape.height

                images.append({
                    "slide": slide_num,
                    "content_type": content_type,
                    "size_bytes": len(blob),
                    "width_emu": width,
                    "height_emu": height,
                })
                print(f"  Slide {slide_num}: {content_type} "
                      f"({len(blob):,} bytes, {width}x{height} EMU)")

    print(f"  Total images extracted: {len(images)}")
    return images


# ---------------------------------------------------------------------------
# 3. DOCX Conversion Tests
# ---------------------------------------------------------------------------

def test_docx_markitdown(docx_path: Path) -> str:
    """Test MarkItDown DOCX conversion (spot check)."""
    from markitdown import MarkItDown

    sub_header("MarkItDown DOCX Conversion (Spot Check)")
    md = MarkItDown()
    result = md.convert(str(docx_path))
    text = result.text_content
    print(f"  Characters: {len(text):,}")
    print(f"  Lines: {len(text.splitlines()):,}")
    print(f"  Headings: {count_headings(text)}")
    print(f"  Links: {count_links(text)}")
    print(f"  Table rows: {count_tables(text)}")
    print(f"  Image refs: {count_images_md(text)}")

    checks = {
        "Title present": "Setup Guide" in text or "Azure KB Agent" in text,
        "Heading preserved": "Prerequisites" in text,
        "Bullet points": "Python 3.11" in text or "Azure CLI" in text,
        "Table content": "AI_SERVICES_ENDPOINT" in text or "SEARCH_ENDPOINT" in text,
        "Numbered list": "azd up" in text or "provision" in text.lower(),
    }
    for label, passed in checks.items():
        status = "✅" if passed else "❌"
        print(f"  {status} {label}")

    # Hyperlink check
    md_links = extract_links(text)
    bare_urls = extract_urls(text)
    print(f"\n  Hyperlinks (Markdown [text](url)): {len(md_links)}")
    for link_text, link_url in md_links[:5]:
        print(f"    [{link_text}]({link_url})")
    print(f"  Bare URLs found: {len(bare_urls)}")
    for url in bare_urls[:5]:
        print(f"    {url}")
    total_urls = len(md_links) + len(bare_urls)
    if total_urls > 0:
        print(f"  ✅ Hyperlinks detected ({total_urls} total)")
    else:
        print(f"  ⚠️  No hyperlinks detected in DOCX conversion output")

    # Long table check
    table_rows = count_tables(text)
    print(f"\n  Table rows in output: {table_rows}")
    if table_rows >= 30:
        print(f"  ✅ Long table (30+ rows) preserved as single table")
    elif table_rows >= 20:
        print(f"  ⚠️  Long table partially extracted ({table_rows} rows)")

    # Image check — all 3 types expected
    images = count_images_md(text)
    print(f"\n  Image refs in output: {images}")
    if images >= 3:
        print(f"  ✅ All 3 image types present in DOCX output")
    elif images > 0:
        print(f"  ⚠️  Only {images} image ref(s) — expected 3")
    else:
        print(f"  ❌ No image refs in DOCX output")

    return text


# ---------------------------------------------------------------------------
# 4. Save outputs and comparison
# ---------------------------------------------------------------------------

def save_output(name: str, text: str) -> Path:
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUTPUT_DIR / f"{name}.md"
    path.write_text(text, encoding="utf-8")
    return path


def save_images(name: str, images: list[dict]) -> list[Path]:
    img_dir = OUTPUT_DIR / f"{name}-images"
    img_dir.mkdir(parents=True, exist_ok=True)
    paths = []
    for img in images:
        if "bytes" in img:
            ext = img.get("ext", "png")
            fname = f"page{img['page']}_img{img['index']}.{ext}"
            path = img_dir / fname
            path.write_bytes(img["bytes"])
            paths.append(path)
    return paths


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    # Ensure samples exist
    if not SAMPLES_DIR.exists():
        print("Samples not found — creating them first...")
        from create_samples import main as create_main
        create_main()

    pdf_path = SAMPLES_DIR / "sample-article.pdf"
    pptx_path = SAMPLES_DIR / "sample-presentation.pptx"
    docx_path = SAMPLES_DIR / "sample-document.docx"

    for p in [pdf_path, pptx_path, docx_path]:
        if not p.exists():
            print(f"ERROR: {p} not found. Run create_samples.py first.")
            sys.exit(1)

    # ===================================================================
    # PDF TESTS (Story 1)
    # ===================================================================
    section_header("STORY 1: PDF CONVERSION & IMAGE EXTRACTION")

    pdf_md = test_pdf_markitdown(pdf_path)
    save_output("pdf-markitdown", pdf_md)

    pdf_md_images = test_pdf_image_extraction_markitdown(pdf_path)

    pdf_pymupdf_images = test_pdf_image_extraction_pymupdf(pdf_path)
    if pdf_pymupdf_images:
        saved = save_images("pdf-pymupdf", pdf_pymupdf_images)
        print(f"  Saved {len(saved)} extracted images to {OUTPUT_DIR / 'pdf-pymupdf-images'}")

    # PDF image extraction comparison
    sub_header("PDF Image Extraction Comparison")
    md_image_count = count_images_md(pdf_md)
    pymupdf_image_count = len(pdf_pymupdf_images)
    print(f"  MarkItDown image refs in Markdown: {md_image_count}")
    print(f"  PyMuPDF extracted images: {pymupdf_image_count}")
    if pymupdf_image_count > 0:
        for img in pdf_pymupdf_images:
            print(f"    - {img['width']}x{img['height']} {img['ext']} "
                  f"({img['size_bytes']:,} bytes)")
        print(f"  ✅ PyMuPDF successfully extracts embedded PDF images")
    if md_image_count == 0:
        print(f"  ⚠️  MarkItDown does not create image references for embedded PDF images")
        print(f"  → PyMuPDF is the recommended approach for PDF image extraction")

    # ===================================================================
    # PPTX TESTS (Story 2)
    # ===================================================================
    section_header("STORY 2: PPTX CONVERSION & SPEAKER NOTES")

    pptx_md = test_pptx_markitdown(pptx_path)
    save_output("pptx-markitdown", pptx_md)

    pptx_notes = test_pptx_speaker_notes_native(pptx_path)

    pptx_images = test_pptx_image_extraction(pptx_path)

    # Compare speaker notes: MarkItDown vs python-pptx
    sub_header("Speaker Notes Comparison")
    notes_in_md = any(
        kw.lower() in pptx_md.lower()
        for kw in ["managed identity", "DefaultAzureCredential", "data flow"]
    )
    notes_via_pptx = any(s["notes_chars"] > 0 for s in pptx_notes)
    print(f"  MarkItDown includes speaker notes: {'✅ Yes' if notes_in_md else '❌ No'}")
    print(f"  python-pptx extracts speaker notes: {'✅ Yes' if notes_via_pptx else '❌ No'}")
    if not notes_in_md and notes_via_pptx:
        print(f"  → python-pptx supplementation needed for speaker notes")

    # ===================================================================
    # DOCX SPOT CHECK (Story 2 addendum)
    # ===================================================================
    section_header("DOCX SPOT CHECK")

    docx_md = test_docx_markitdown(docx_path)
    save_output("docx-markitdown", docx_md)

    # ===================================================================
    # SUMMARY TABLE
    # ===================================================================
    section_header("SUMMARY: CONVERSION QUALITY MATRIX")

    print(f"\n{'Format':<10} {'Chars':>8} {'Lines':>8} {'Headings':>10} "
          f"{'Links':>8} {'Tables':>8} {'Images':>8}")
    print("-" * 62)
    for label, text in [
        ("PDF", pdf_md),
        ("PPTX", pptx_md),
        ("DOCX", docx_md),
    ]:
        print(f"{label:<10} {len(text):>8,} {len(text.splitlines()):>8,} "
              f"{count_headings(text):>10} {count_links(text):>8} "
              f"{count_tables(text):>8} {count_images_md(text):>8}")

    # ===================================================================
    # FINDINGS SUMMARY
    # ===================================================================
    section_header("KEY FINDINGS")

    findings = []

    # PDF findings
    if count_headings(pdf_md) > 0:
        findings.append(("PDF text extraction", "✅", "MarkItDown extracts text and headings from PDF"))
    else:
        findings.append(("PDF text extraction", "⚠️", "MarkItDown PDF text extraction limited"))

    if count_tables(pdf_md) > 0:
        findings.append(("PDF table extraction", "✅", "Tables preserved in MarkItDown PDF output"))
    else:
        findings.append(("PDF table extraction", "❌", "Tables NOT preserved in MarkItDown PDF output"))

    if pymupdf_image_count > 0:
        findings.append(("PDF image extraction", "✅",
                        f"PyMuPDF extracts {pymupdf_image_count} images at original resolution"))
    else:
        findings.append(("PDF image extraction", "⚠️", "No images extracted (may be test PDF limitation)"))

    # PPTX findings
    if "Architecture" in pptx_md or "System Components" in pptx_md:
        findings.append(("PPTX text extraction", "✅", "Slide content extracted by MarkItDown"))
    else:
        findings.append(("PPTX text extraction", "❌", "Slide content NOT extracted"))

    if notes_in_md:
        findings.append(("PPTX speaker notes", "✅", "MarkItDown includes speaker notes in output"))
    else:
        findings.append(("PPTX speaker notes", "❌", "Speaker notes NOT in MarkItDown output — python-pptx needed"))

    # DOCX findings
    if count_headings(docx_md) > 0 and count_tables(docx_md) > 0:
        findings.append(("DOCX conversion", "✅", "Headings, tables, and lists preserved"))
    else:
        findings.append(("DOCX conversion", "⚠️", "Some content not preserved"))

    # Cross-format: hyperlink extraction
    for fmt, md_text in [("PDF", pdf_md), ("PPTX", pptx_md), ("DOCX", docx_md)]:
        md_links = extract_links(md_text)
        bare_urls = extract_urls(md_text)
        total = len(md_links) + len(bare_urls)
        if len(md_links) > 0:
            findings.append((f"{fmt} hyperlinks", "✅",
                            f"{len(md_links)} Markdown links extracted — no post-processing needed"))
        elif len(bare_urls) > 0:
            findings.append((f"{fmt} hyperlinks", "⚠️",
                            f"{len(bare_urls)} bare URLs found but not as [text](url) — needs post-processing"))
        else:
            findings.append((f"{fmt} hyperlinks", "❌",
                            "No hyperlinks detected — links lost during conversion"))

    # Cross-format: long table handling
    for fmt, md_text in [("PDF", pdf_md), ("PPTX", pptx_md), ("DOCX", docx_md)]:
        rows = count_tables(md_text)
        if rows >= 30:
            findings.append((f"{fmt} long table", "✅",
                            f"Long table preserved as single table ({rows} rows)"))
        elif rows >= 20:
            findings.append((f"{fmt} long table", "⚠️",
                            f"Long table partially extracted ({rows} rows — may be split)"))
        else:
            findings.append((f"{fmt} long table", "⚠️",
                            f"Long table may be split across pages ({rows} table rows total)"))

    for label, status, detail in findings:
        print(f"  {status} {label}: {detail}")

    # ===================================================================
    # RECOMMENDATION
    # ===================================================================
    section_header("RECOMMENDATION")

    critical_failures = sum(1 for _, s, _ in findings if s == "❌")
    warnings = sum(1 for _, s, _ in findings if s == "⚠️")
    successes = sum(1 for _, s, _ in findings if s == "✅")

    print(f"\n  Results: {successes} ✅ pass, {warnings} ⚠️ warnings, {critical_failures} ❌ failures")

    if critical_failures == 0:
        print("\n  🟢 GO — MarkItDown is viable for PDF/PPTX conversion with supplementation.")
    elif critical_failures <= 2:
        print("\n  🟡 CONDITIONAL GO — MarkItDown needs supplementation for identified gaps.")
    else:
        print("\n  🔴 NO-GO — MarkItDown has too many gaps for PDF/PPTX conversion.")

    print("\n  Mitigations:")
    if not notes_in_md:
        print("    - Use python-pptx to extract speaker notes and append to MarkItDown output")
    if md_image_count == 0 and pymupdf_image_count > 0:
        print("    - Use PyMuPDF (fitz) for PDF image extraction instead of MarkItDown")
    # Hyperlink mitigation
    for fmt, md_text in [("PDF", pdf_md), ("PPTX", pptx_md), ("DOCX", docx_md)]:
        md_links = extract_links(md_text)
        bare_urls = extract_urls(md_text)
        if len(md_links) == 0 and len(bare_urls) > 0:
            print(f"    - {fmt}: Bare URLs present but not Markdown-formatted — needs post-processing to wrap in [text](url)")
        elif len(md_links) == 0 and len(bare_urls) == 0:
            print(f"    - {fmt}: Hyperlinks lost — extract from source format before/after conversion")
    print("    - DOCX: MarkItDown's built-in support is sufficient (no supplementation needed)")

    print()


if __name__ == "__main__":
    main()
