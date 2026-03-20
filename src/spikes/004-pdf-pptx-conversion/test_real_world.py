#!/usr/bin/env python3
"""Test MarkItDown conversion quality on real-world public documents.

Runs conversion analysis on any documents found in samples/real-world/.
Complements the synthetic tests in run.py with real-world complexity:
multi-column layouts, real photographs, dense hyperlinks, and varied formatting.

Usage:
    cd src/functions && uv run python ../spikes/004-pdf-pptx-conversion/test_real_world.py

If no real-world documents are found, run download_real_world.py first.
"""

from __future__ import annotations

import re
import sys
from pathlib import Path

SPIKE_DIR = Path(__file__).resolve().parent
SAMPLES_DIR = SPIKE_DIR / "samples" / "real-world"
OUTPUT_DIR = SPIKE_DIR / "output" / "real-world"


# ---------------------------------------------------------------------------
# Helpers (shared with run.py)
# ---------------------------------------------------------------------------

def count_headings(md: str) -> int:
    return len(re.findall(r"^#{1,6}\s+", md, re.MULTILINE))


def extract_md_links(md: str) -> list[tuple[str, str]]:
    """Extract [text](url) hyperlinks (excluding image refs)."""
    return re.findall(r"(?<!!)\[([^\]]+)\]\((https?://[^)]+)\)", md)


def extract_bare_urls(md: str) -> list[str]:
    return re.findall(r"https?://[^\s\)>\"]+", md)


def count_table_rows(md: str) -> int:
    return len(re.findall(r"^\|.+\|$", md, re.MULTILINE))


def count_image_refs(md: str) -> int:
    return len(re.findall(r"!\[([^\]]*)\]\(([^)]+)\)", md))


def section(title: str) -> None:
    print(f"\n{'=' * 70}")
    print(f"  {title}")
    print(f"{'=' * 70}")


def sub(title: str) -> None:
    print(f"\n--- {title} ---")


# ---------------------------------------------------------------------------
# PDF analysis
# ---------------------------------------------------------------------------

def analyze_pdf(pdf_path: Path) -> dict:
    """Run MarkItDown + PyMuPDF analysis on a real-world PDF."""
    from markitdown import MarkItDown

    sub(f"MarkItDown: {pdf_path.name}")
    md = MarkItDown()
    result = md.convert(str(pdf_path))
    text = result.text_content

    stats = {
        "file": pdf_path.name,
        "format": "PDF",
        "size_bytes": pdf_path.stat().st_size,
        "chars": len(text),
        "lines": len(text.splitlines()),
        "headings": count_headings(text),
        "md_links": len(extract_md_links(text)),
        "bare_urls": len(extract_bare_urls(text)),
        "table_rows": count_table_rows(text),
        "image_refs": count_image_refs(text),
    }
    for k, v in stats.items():
        if k not in ("file", "format"):
            print(f"  {k}: {v:,}" if isinstance(v, int) else f"  {k}: {v}")

    # PyMuPDF image + link extraction
    try:
        import fitz
        doc = fitz.open(str(pdf_path))
        total_images = 0
        total_links = 0
        for page_num in range(len(doc)):
            page = doc[page_num]
            total_images += len(page.get_images(full=True))
            total_links += len(page.get_links())
        doc.close()
        stats["pymupdf_images"] = total_images
        stats["pymupdf_links"] = total_links
        print(f"  pymupdf_images: {total_images}")
        print(f"  pymupdf_links: {total_links}")
    except ImportError:
        print("  ⚠️  PyMuPDF not installed — skipping image/link count")

    # Save converted output
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / f"{pdf_path.stem}-markitdown.md"
    out_path.write_text(text, encoding="utf-8")
    print(f"  output: {out_path}")

    # Key quality checks
    print()
    if stats["md_links"] > 0:
        print(f"  ✅ Hyperlinks: {stats['md_links']} Markdown links extracted")
    elif stats["bare_urls"] > 0:
        print(f"  ⚠️  Hyperlinks: {stats['bare_urls']} bare URLs found (not [text](url))")
    else:
        total_source = stats.get("pymupdf_links", 0)
        print(f"  ❌ Hyperlinks: 0 extracted (source has {total_source} links)")

    if stats["headings"] > 0:
        print(f"  ✅ Headings: {stats['headings']} with # markers")
    else:
        print(f"  ⚠️  Headings: none with # markers (headings as plain text)")

    if stats["table_rows"] > 0:
        print(f"  ✅ Tables: {stats['table_rows']} pipe-delimited rows")
    else:
        print(f"  ⚠️  Tables: none detected in Markdown output")

    pymupdf_imgs = stats.get("pymupdf_images", 0)
    if stats["image_refs"] > 0:
        print(f"  ✅ Images: {stats['image_refs']} refs in Markdown")
    elif pymupdf_imgs > 0:
        print(f"  ❌ Images: 0 in Markdown (PyMuPDF found {pymupdf_imgs} in source)")
    else:
        print(f"  — Images: none in source document")

    return stats


# ---------------------------------------------------------------------------
# PPTX analysis
# ---------------------------------------------------------------------------

def analyze_pptx(pptx_path: Path) -> dict:
    """Run MarkItDown analysis on a real-world PPTX."""
    from markitdown import MarkItDown

    sub(f"MarkItDown: {pptx_path.name}")
    md = MarkItDown()
    result = md.convert(str(pptx_path))
    text = result.text_content

    stats = {
        "file": pptx_path.name,
        "format": "PPTX",
        "size_bytes": pptx_path.stat().st_size,
        "chars": len(text),
        "lines": len(text.splitlines()),
        "headings": count_headings(text),
        "md_links": len(extract_md_links(text)),
        "bare_urls": len(extract_bare_urls(text)),
        "table_rows": count_table_rows(text),
        "image_refs": count_image_refs(text),
    }
    for k, v in stats.items():
        if k not in ("file", "format"):
            print(f"  {k}: {v:,}" if isinstance(v, int) else f"  {k}: {v}")

    # Check for speaker notes
    notes_pattern = r"###\s*Notes:"
    notes_count = len(re.findall(notes_pattern, text))
    stats["speaker_notes_sections"] = notes_count
    print(f"  speaker_notes_sections: {notes_count}")

    # python-pptx cross-check
    try:
        from pptx import Presentation
        prs = Presentation(str(pptx_path))
        slides_with_notes = sum(
            1 for s in prs.slides
            if s.has_notes_slide and s.notes_slide.notes_text_frame.text.strip()
        )
        total_slides = len(prs.slides)
        stats["total_slides"] = total_slides
        stats["slides_with_notes"] = slides_with_notes
        print(f"  total_slides: {total_slides}")
        print(f"  slides_with_notes: {slides_with_notes}")
    except Exception:
        pass

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / f"{pptx_path.stem}-markitdown.md"
    out_path.write_text(text, encoding="utf-8")
    print(f"  output: {out_path}")

    print()
    if notes_count > 0:
        print(f"  ✅ Speaker notes: {notes_count} ### Notes: sections found")
    else:
        slides_notes = stats.get("slides_with_notes", 0)
        if slides_notes > 0:
            print(f"  ❌ Speaker notes: {slides_notes} slides have notes but NOT in MarkItDown output")
        else:
            print(f"  — Speaker notes: none in source PPTX")

    if stats["image_refs"] > 0:
        print(f"  ✅ Images: {stats['image_refs']} refs in Markdown")

    if stats["md_links"] > 0:
        print(f"  ✅ Hyperlinks: {stats['md_links']} Markdown links extracted")
    elif stats["bare_urls"] > 0:
        print(f"  ⚠️  Hyperlinks: {stats['bare_urls']} bare URLs (not [text](url))")
    else:
        print(f"  ⚠️  Hyperlinks: none detected")

    return stats


# ---------------------------------------------------------------------------
# DOCX analysis
# ---------------------------------------------------------------------------

def analyze_docx(docx_path: Path) -> dict:
    """Run MarkItDown analysis on a real-world DOCX."""
    from markitdown import MarkItDown

    sub(f"MarkItDown: {docx_path.name}")
    md = MarkItDown()
    result = md.convert(str(docx_path))
    text = result.text_content

    stats = {
        "file": docx_path.name,
        "format": "DOCX",
        "size_bytes": docx_path.stat().st_size,
        "chars": len(text),
        "lines": len(text.splitlines()),
        "headings": count_headings(text),
        "md_links": len(extract_md_links(text)),
        "bare_urls": len(extract_bare_urls(text)),
        "table_rows": count_table_rows(text),
        "image_refs": count_image_refs(text),
    }
    for k, v in stats.items():
        if k not in ("file", "format"):
            print(f"  {k}: {v:,}" if isinstance(v, int) else f"  {k}: {v}")

    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    out_path = OUTPUT_DIR / f"{docx_path.stem}-markitdown.md"
    out_path.write_text(text, encoding="utf-8")
    print(f"  output: {out_path}")

    print()
    if stats["md_links"] > 0:
        print(f"  ✅ Hyperlinks: {stats['md_links']} Markdown links extracted")
    elif stats["bare_urls"] > 0:
        print(f"  ⚠️  Hyperlinks: {stats['bare_urls']} bare URLs (not [text](url))")
    else:
        print(f"  ⚠️  Hyperlinks: none detected")

    if stats["headings"] > 0:
        print(f"  ✅ Headings: {stats['headings']} with # markers")
    if stats["table_rows"] > 0:
        print(f"  ✅ Tables: {stats['table_rows']} pipe-delimited rows")
    if stats["image_refs"] > 0:
        print(f"  ✅ Images: {stats['image_refs']} refs in Markdown")

    return stats


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main() -> None:
    if not SAMPLES_DIR.exists():
        print(f"No real-world samples found at {SAMPLES_DIR}")
        print("Run download_real_world.py first.")
        sys.exit(1)

    # Discover files
    pdfs = sorted(SAMPLES_DIR.glob("*.pdf"))
    pptxs = sorted(SAMPLES_DIR.glob("*.pptx"))
    docxs = sorted(SAMPLES_DIR.glob("*.docx"))

    total = len(pdfs) + len(pptxs) + len(docxs)
    if total == 0:
        print(f"No documents found in {SAMPLES_DIR}")
        print("Run download_real_world.py first.")
        sys.exit(1)

    section("REAL-WORLD DOCUMENT CONVERSION TEST")
    print(f"\n  Found: {len(pdfs)} PDF, {len(pptxs)} PPTX, {len(docxs)} DOCX")

    all_stats = []

    # --- PDFs ---
    if pdfs:
        section("PDF CONVERSION (Real-World)")
        for pdf in pdfs:
            stats = analyze_pdf(pdf)
            all_stats.append(stats)

    # --- PPTXs ---
    if pptxs:
        section("PPTX CONVERSION (Real-World)")
        for pptx in pptxs:
            stats = analyze_pptx(pptx)
            all_stats.append(stats)

    # --- DOCXs ---
    if docxs:
        section("DOCX CONVERSION (Real-World)")
        for docx in docxs:
            stats = analyze_docx(docx)
            all_stats.append(stats)

    # --- Summary ---
    section("REAL-WORLD CONVERSION QUALITY MATRIX")

    print(f"\n  {'File':<35s} {'Format':<6s} {'Size':>10s} {'Chars':>8s} "
          f"{'Lines':>7s} {'Heads':>6s} {'Links':>6s} {'Tables':>7s} {'Imgs':>5s}")
    print("  " + "-" * 96)
    for s in all_stats:
        print(f"  {s['file']:<35s} {s['format']:<6s} "
              f"{s['size_bytes']:>10,} {s['chars']:>8,} "
              f"{s['lines']:>7,} {s['headings']:>6} "
              f"{s['md_links']:>6} {s['table_rows']:>7} "
              f"{s['image_refs']:>5}")

    # --- Key Findings ---
    section("REAL-WORLD KEY FINDINGS")

    for s in all_stats:
        print(f"\n  📄 {s['file']} ({s['format']})")

        # Hyperlink gap
        if s["format"] == "PDF":
            pymupdf_links = s.get("pymupdf_links", 0)
            if s["md_links"] == 0 and pymupdf_links > 0:
                print(f"    ❌ Hyperlinks: {pymupdf_links} in source → 0 in output "
                      f"(CONFIRMED: PDF hyperlinks lost)")
            elif s["md_links"] > 0:
                print(f"    ✅ Hyperlinks: {s['md_links']} extracted")

            pymupdf_imgs = s.get("pymupdf_images", 0)
            if s["image_refs"] == 0 and pymupdf_imgs > 0:
                print(f"    ❌ Images: {pymupdf_imgs} in source → 0 in output "
                      f"(CONFIRMED: PDF images need PyMuPDF)")
        elif s["format"] == "PPTX":
            if s["md_links"] == 0:
                print(f"    ❌ Hyperlinks: lost (bare URLs: {s['bare_urls']})")
            notes = s.get("speaker_notes_sections", 0)
            src_notes = s.get("slides_with_notes", 0)
            if notes > 0:
                print(f"    ✅ Speaker notes: {notes} sections")
            elif src_notes > 0:
                print(f"    ❌ Speaker notes: {src_notes} in source → 0 in output")
        elif s["format"] == "DOCX":
            if s["md_links"] > 0:
                print(f"    ✅ Hyperlinks: {s['md_links']} preserved")
            else:
                print(f"    ⚠️  Hyperlinks: {s['bare_urls']} bare URLs only")

    print()


if __name__ == "__main__":
    main()
